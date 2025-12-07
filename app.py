import gradio as gr
from openai import OpenAI
from pypdf import PdfReader
import docx
from pptx import Presentation
import os

# Use the OPENAI_API_KEY secret from Hugging Face
client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

# -------------------------------------------------------------------
# Helper: extract text from PDFs, Word docs, PPTs
# -------------------------------------------------------------------

def extract_text(file):
    # Some uploads may come as tempfile objects; use file.name safely
    name = getattr(file, "name", "")
    if name.endswith(".pdf"):
        reader = PdfReader(file)
        return "\n".join([page.extract_text() or "" for page in reader.pages])

    if name.endswith(".docx"):
        d = docx.Document(file)
        return "\n".join([p.text for p in d.paragraphs])

    if name.endswith(".pptx"):
        pres = Presentation(file)
        text = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    # fallback: plain text
    try:
        return file.read().decode("utf-8", errors="ignore")
    except Exception:
        return ""

# -------------------------------------------------------------------
# AI Generator
# -------------------------------------------------------------------

def generate_all(files, subject_name, exam_instructions):
    # If no subject name, use a default instead of crashing UI
    if not subject_name:
        subject_name = "My Course"

    # If no files, show a friendly message
    if not files:
        msg = "Please upload at least one syllabus / slide / PDF to generate a study pack."
        return msg, "", "", ""

    # combine text from all uploaded files
    full_corpus = ""
    for f in files:
        full_corpus += extract_text(f) + "\n\n"

    prompt = f"""
You are an AI tutor creating a full study system for the subject: {subject_name}.

Here is the course material (syllabi, slides, notes, etc.):
{full_corpus}

Please generate ALL of the following in one message, separated clearly with headings:

1. PRAACHI-GUIDE TEXTBOOK:
- Written in Praachi’s casual academic voice:
  - clear, simple explanations
  - conversational but not sloppy
  - uses intuition and examples
  - points out common mistakes
- Organized into sections/chapters
- Include short knowledge checks every few paragraphs.

2. FLASHCARDS:
Provide 15–25 flashcards in this exact format:
Q: <front of card>
A: <back of card>

3. QUESTION BANK:
Provide 10 multiple-choice questions + 10 short-answer questions.
For each MCQ, include:
- The correct answer
- Why each wrong option is wrong
- A short reference label like [See: Topic X in Textbook]

4. EXAM TEMPLATE:
Generate an exam using these optional instructions:
\"\"\"{exam_instructions}\"\"\"

Include:
- 10 multiple-choice questions
- 5 short-answer questions
- 1 longer applied / case-style question.
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
    )

    output = response.choices[0].message.content

    # We split the output by known markers. This doesn't have to be perfect;
    # it just needs to roughly separate sections.
    # If a marker is missing, fall back gracefully.
    textbook = ""
    flashcards = ""
    question_bank = ""
    exam = ""

    # Normalize headings to make splitting easier
    normalized = output.replace("Flashcards", "FLASHCARDS").replace("Question Bank", "QUESTION BANK").replace("Exam Template", "EXAM TEMPLATE")

    if "FLASHCARDS" in normalized:
        parts = normalized.split("FLASHCARDS", 1)
        textbook = parts[0].replace("PRAACHI-GUIDE TEXTBOOK", "").strip()
        rest = parts[1]
    else:
        textbook = normalized.strip()
        rest = ""

    if "QUESTION BANK" in rest:
        parts2 = rest.split("QUESTION BANK", 1)
        flashcards = parts2[0].strip()
        rest2 = parts2[1]
    else:
        flashcards = rest.strip()
        rest2 = ""

    if "EXAM TEMPLATE" in rest2:
        parts3 = rest2.split("EXAM TEMPLATE", 1)
        question_bank = parts3[0].strip()
        exam = parts3[1].strip()
    else:
        question_bank = rest2.strip()

    return textbook, flashcards, question_bank, exam

# -------------------------------------------------------------------
# Exam feedback generator
# -------------------------------------------------------------------

def grade_exam(user_answers, exam_text):
    if not user_answers:
        return "Please paste your exam answers before submitting."

    prompt = f"""
You are grading the following exam.

Exam text:
{exam_text}

Student answers:
{user_answers}

Please give:
1. A score out of 100.
2. What they understood well.
3. What they misunderstood.
4. Specific topics and textbook sections they should review (refer to themes, not exact page numbers).
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
    )

    return response.choices[0].message.content

# -------------------------------------------------------------------
# Gradio UI
# -------------------------------------------------------------------

with gr.Blocks() as demo:
    gr.Markdown("# 📚 AI Study Tool (Praachi-Style Edition)")
    gr.Markdown("Upload your course materials and auto-generate a textbook, flashcards, a question bank, and a full exam.")

    subject_name = gr.Textbox(label="Subject Name")
    exam_instr = gr.Textbox(label="Custom Exam Instructions (optional)", lines=6)
    file_upload = gr.File(label="Upload syllabus, slides, PDFs, docs, etc.", file_count="multiple")

    generate_btn = gr.Button("Generate Study Pack")

    with gr.Tab("Textbook"):
        textbook_box = gr.Markdown()

    with gr.Tab("Flashcards"):
        flashcard_box = gr.Markdown()

    with gr.Tab("Question Bank"):
        question_box = gr.Markdown()

    with gr.Tab("Exam"):
        exam_box = gr.Markdown()
        user_answers = gr.Textbox(label="Paste your exam answers here", lines=8)
        grade_btn = gr.Button("Submit Exam for Feedback")
        grade_output = gr.Markdown()

    generate_all_outputs = [textbook_box, flashcard_box, question_box, exam_box]

    generate_btn.click(
        generate_all,
        inputs=[file_upload, subject_name, exam_instr],
        outputs=generate_all_outputs,
    )

    grade_btn.click(
        grade_exam,
        inputs=[user_answers, exam_box],
        outputs=grade_output,
    )

demo.launch()
